from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import re

class PowerPointFormatter:
    def __init__(self):
        self.prs = Presentation()
        
    def create_title_slide(self, title, subtitle=None):
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Set smaller font size for title
        title_frame = title_shape.text_frame
        title_frame.clear()
        title_para = title_frame.add_paragraph()
        title_para.text = title
        title_para.font.size = Pt(32)  # Smaller title font size
        
        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            
            # Set smaller font size for subtitle
            subtitle_frame = subtitle_shape.text_frame
            subtitle_frame.clear()
            subtitle_para = subtitle_frame.add_paragraph()
            subtitle_para.text = subtitle
            subtitle_para.font.size = Pt(18)  # Smaller subtitle font size
            
        return slide
    
    def create_content_slide(self, title, content_points):
        slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Set smaller font size for slide title
        title_frame = title_shape.text_frame
        title_frame.clear()
        title_para = title_frame.add_paragraph()
        title_para.text = title
        title_para.font.size = Pt(24)  # Smaller slide title font size
        
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.text = ""  # Clear default text
        
        for point in content_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.space_after = Pt(12)
            p.font.size = Pt(14)  # Smaller content font size
            
        return slide
    
    def save_presentation(self, filename, output_dir="output"):
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
        
        filepath = os.path.join(output_dir, f"{filename}.pptx")
        self.prs.save(filepath)
        return filepath
    
    def format_content_from_outline(self, outline):
        # Parse the structured outline and create slides
        # This is a simplified implementation - you might need to adjust based on your actual outline format
        text = re.sub(r"\*\*", "", outline).strip()
        
        lines = text.split("\n")
        
        slides_data = []
        current_slide_title = None
        current_points = []

        for line in lines:
            line = line.strip()
            
            # Detect slide titles like "Slide 1: Title" or "Slide 2: ..."
            if re.match(r"^Slide\s*\d+:", line, re.IGNORECASE):
                # Save previous slide if exists
                if current_slide_title:
                    slides_data.append({
                        'title': current_slide_title,
                        'points': current_points
                    })
                # Start new slide
                current_slide_title = line.split(":", 1)[1].strip()
                current_points = []
            
            # Detect bullets (even nested ones)
            elif line.startswith("-"):
                # Remove leading '-' and extra spaces
                point = line.lstrip("-").strip()
                if point:
                    current_points.append(point)
        
        # Add last slide
        if current_slide_title and current_points:
            slides_data.append({
                'title': current_slide_title,
                'points': current_points
            })
        
        # Generate slides
        for i, slide_data in enumerate(slides_data):
            if i == 0:
                subtitle = slide_data['points'][0] if slide_data['points'] else ""
                self.create_title_slide(slide_data['title'], subtitle)
            else:
                self.create_content_slide(slide_data['title'], slide_data['points'])
        
        return self 
    
    